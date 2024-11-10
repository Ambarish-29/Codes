import pandas as pd
from pptx import Presentation
from pptx.util import Inches

# Function to format date with correct suffix
def add_date_suffix(day):
    if 11 <= day <= 13:
        suffix = 'th'
    else:
        suffix = {1: 'st', 2: 'nd', 3: 'rd'}.get(day % 10, 'th')
    return suffix

def format_date(date):
    day = date.day
    suffix = add_date_suffix(day)
    return f"{day}{suffix} {date.strftime('%b')}"

date_filter = True
if date_filter:
    start_date = "2023-01-01"
    end_date = "2025-05-07"

df = pd.read_excel("weekly.xlsx")

# Convert Date column to datetime
df['Date'] = pd.to_datetime(df['Date'], format='%m/%d/%Y')

df = df[(df['Date'] >= start_date) & (df['Date'] <= end_date)]

# Concatenate Date and campaign with unique dates and proper sorting
df['campaign_date'] = df.groupby(['product', 'campaign'])['Date'].transform(
    lambda x: ', '.join([format_date(date) for date in sorted(set(x), key=lambda date: date.strftime('%Y-%m-%d'))]))

# Create a cross-tabulation
pivot_df = pd.crosstab(index=[df['product'], df['campaign'], df['campaign_date']], columns=df['channel'])

# Calculate Grand Total
grand_total = pivot_df.sum()
grand_total.name = ('Grand Total', '', '')
pivot_df.loc[('Grand Total', '', '')] = grand_total

# Convert float values to integers
pivot_df = pivot_df.astype(int)

# Calculate product totals
product_total = pivot_df.groupby(level='product').sum().reset_index()
product_total['campaign'] = ''
product_total['campaign_date'] = ''
product_total.set_index(['product', 'campaign', 'campaign_date'], inplace=True)

# Remove the second "Grand Total" row
pivot_df = pivot_df.loc[pivot_df.index != ('Grand Total', '', '')]

# Concatenate product totals with the original pivot table
pivot_df = pd.concat([product_total, pivot_df])

# Sort rows based on alphabetical order of product names
pivot_df = pivot_df.sort_index(level='product')

# Calculate total of email, push, sms for each row
pivot_df['Channel total'] = pivot_df.sum(axis=1)

# Calculate the sum of "size" according to the group
size_sum = df.groupby(['product', 'campaign', 'campaign_date'])['size'].sum()
prod_size_sum = df.groupby(['product'])['size'].sum()
prod_size_sum.index = pd.MultiIndex.from_tuples([(product, '', '') for product in prod_size_sum.index], names=['product', 'campaign', 'campaign_date'])
grand_total_size = size_sum.sum()
grand_total_df = pd.DataFrame([[grand_total_size]], columns=['size'], index=pd.MultiIndex.from_tuples([('Grand Total', '', '')], names=['product', 'campaign', 'campaign_date']))
test_df = pd.concat([size_sum, prod_size_sum, grand_total_df])
test_df = test_df.sort_index(level='product')

test_df = test_df.reindex(pivot_df.index, fill_value=0)  # Align the index
pivot_df['sum_of_size'] = test_df  # Add the calculated sums as a new column

pivot_df = pivot_df.astype(int)
print(pivot_df)

# Create a PowerPoint presentation object
prs = Presentation()

# Add a slide with a title and content layout
slide_layout = prs.slide_layouts[5]  # Use the title and content layout
slide = prs.slides.add_slide(slide_layout)

# Set the title for the slide
title = slide.shapes.title
title.text = "Pivot Table"

# Add the pivot table as a table on the slide
rows, cols = pivot_df.shape
table = slide.shapes.add_table(rows + 1, cols + 3, Inches(1), Inches(2), Inches(8), Inches(4)).table

# Add headers
for col_idx, col_name in enumerate(['product', 'campaign', 'campaign_date'] + list(pivot_df.columns)):
    table.cell(0, col_idx).text = str(col_name)

# Add data to the table, maintaining the structure
prev_product = None
prev_campaign = None
for row_idx, (index, row) in enumerate(pivot_df.iterrows(), start=1):
    product, campaign, campaign_date = index

    if product == prev_product:
        table.cell(row_idx, 0).text = ''
    else:
        table.cell(row_idx, 0).text = product
        prev_product = product

    if campaign == prev_campaign and campaign_date == '':
        table.cell(row_idx, 1).text = ''
        table.cell(row_idx, 2).text = ''
    else:
        table.cell(row_idx, 1).text = campaign
        table.cell(row_idx, 2).text = campaign_date
        prev_campaign = campaign

    for col_idx, value in enumerate(row, start=3):
        table.cell(row_idx, col_idx).text = str(value)

# Save the presentation
pptx_file = 'pivot_table_presentation.pptx'
prs.save(pptx_file)

print(f'Presentation saved as {pptx_file}')

pivot_df.to_excel('weekly_dashboard.xlsx')

email_df = pd.read_excel("weekly_email.xlsx")

email_pivot = email_df.pivot_table(index=['type', 'product', 'camp'], values=['size', 'open', 'click'], aggfunc='sum')

print(email_pivot)
